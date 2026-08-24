# engines/export-engine/pptx_format.py
#
# WHAT: Deterministic PowerPoint renderer for Journeys — one deck per
#       journey: a title slide plus one slide per Card (content + quiz).
# WHY:  MASTER_STORY promises PPTX among the export formats (P4.3).
#       Core properties are pinned so identical input yields byte-identical
#       output (P4.4 exit criterion E4).
# BREAKS IF DELETED: Journeys can no longer be exported to PowerPoint.

from __future__ import annotations

import logging
from datetime import datetime
from io import BytesIO
from typing import Any

logger = logging.getLogger(__name__)

_EPOCH = datetime(2000, 1, 1)


def export_journey_to_pptx(journey: dict[str, Any]) -> bytes:
    """
    Contract: render a validated Journey dict to PPTX bytes.

    Slide 1 is a title slide (topic + level + card count); every Card gets
    its own slide with content followed by the quiz question, options,
    correct answer, and explanation.

    Raises:
        ValueError: missing topic/cards.
        ImportError: python-pptx not installed.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ImportError(
            "pptx export requires python-pptx: pip install python-pptx"
        ) from exc

    topic = journey.get("topic")
    level = journey.get("level", "beginner")
    cards = journey.get("cards", [])
    if not topic:
        raise ValueError("Journey must have a 'topic' field")
    if not cards:
        raise ValueError("Journey must contain at least one card")

    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_slide.placeholders[1].text = (
        f"{level.title()} level · {len(cards)} cards"
    )

    body_layout = prs.slide_layouts[1]  # title + content
    for i, card in enumerate(cards, start=1):
        slide = prs.slides.add_slide(body_layout)
        slide.shapes.title.text = f"Card {i}: {card.get('title', '')}"
        lines = [
            card.get("content", ""),
            "",
            f"Quiz: {card.get('question', '')}",
        ]
        for j, opt in enumerate(card.get("options", [])):
            lines.append(f"{chr(65 + j)}) {opt}")
        lines += [
            f"Correct: {card.get('correct_option', '')}",
            f"Why: {card.get('explanation', '')}",
        ]
        slide.placeholders[1].text = "\n".join(lines)

    props = prs.core_properties
    props.author = "L&D Command Center"
    props.created = _EPOCH
    props.modified = _EPOCH
    props.title = topic

    buf = BytesIO()
    prs.save(buf)
    logger.info("Rendered PPTX export for %r (%d cards)", topic, len(cards))
    return buf.getvalue()
